"""PowerPoint (.pptx) -> slide text + speaker notes + embedded tables."""
from __future__ import annotations

from pathlib import Path

from ..types import ParsedDoc, Table


def parse(path: str | Path) -> ParsedDoc:
    """Extract text from every shape and speaker notes per slide.

    Each slide is prefixed with a '--- Slide N ---' header so the chunker
    produces slide-scoped prose chunks. Embedded tables become Table objects.
    Returns ParsedDoc with text=joined slides and tables=list of Table.
    """
    from pptx import Presentation

    path = Path(path)
    prs = Presentation(path)

    slide_texts: list[str] = []
    tables: list[Table] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts: list[str] = []

        for shape in slide.shapes:
            # Text frames (title, content, text boxes)
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)

            # Embedded tables
            if shape.has_table:
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                tables.append(Table(name=f"slide_{slide_num}_table", rows=rows))

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Speaker notes]: {notes_text}")

        if parts:
            slide_texts.append(f"--- Slide {slide_num} ---\n" + "\n".join(parts))

    text = "\n\n".join(slide_texts)

    return ParsedDoc(
        text=text,
        tables=tables,
        metadata={
            "filename": path.name,
            "type": "pptx",
            "slides": len(prs.slides),
        },
        source_ref={"filename": path.name, "slide_count": len(prs.slides)},
    )
